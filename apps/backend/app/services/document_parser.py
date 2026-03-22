from dataclasses import dataclass
import logging
from pathlib import Path
import base64

from docx import Document
from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".ppt"}
logging.getLogger("pypdf").setLevel(logging.ERROR)
# add ppt support

@dataclass
class TextSegment:
    index: int
    source_unit: str
    text: str


@dataclass
class ParsedDocument:
    file_path: Path
    doc_type: str
    segments: list[TextSegment]

    @property
    def full_text(self) -> str:
        return "\n".join(seg.text for seg in self.segments if seg.text.strip())

    @property
    def text_chars(self) -> int:
        return len(self.full_text)

    @property
    def segment_count(self) -> int:
        return len(self.segments)


def parse_document(file_path: Path, max_pdf_pages: int = 80) -> ParsedDocument:
    suffix = file_path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        return ParsedDocument(file_path=file_path, doc_type=suffix.lstrip("."), segments=[])

    if suffix in {".txt", ".md"}:
        content = file_path.read_text(encoding="utf-8", errors="ignore")
        segments = [
            TextSegment(index=i, source_unit=f"line_{i+1}", text=line.strip())
            for i, line in enumerate(content.splitlines())
            if line.strip()
        ]
        return ParsedDocument(file_path=file_path, doc_type=suffix.lstrip("."), segments=segments)

    if suffix == ".pdf":
        reader = PdfReader(str(file_path), strict=False)
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:  # noqa: BLE001
                return ParsedDocument(file_path=file_path, doc_type="pdf", segments=[])
        segments: list[TextSegment] = []
        for i, page in enumerate(reader.pages[:max_pdf_pages]):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                segments.append(TextSegment(index=i, source_unit=f"page_{i+1}", text=page_text))
        return ParsedDocument(file_path=file_path, doc_type="pdf", segments=segments)

    if suffix == ".docx":
        doc = Document(str(file_path))
        segments = [
            TextSegment(index=i, source_unit=f"paragraph_{i+1}", text=p.text.strip())
            for i, p in enumerate(doc.paragraphs)
            if p.text.strip()
        ]
        return ParsedDocument(file_path=file_path, doc_type="docx", segments=segments)

    if suffix == ".pptx":
        prs = Presentation(str(file_path))
        segments: list[TextSegment] = []
        for i, slide in enumerate(prs.slides):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_lines.append(text.strip())
            if slide_lines:
                segments.append(
                    TextSegment(
                        index=i,
                        source_unit=f"slide_{i+1}",
                        text="\n".join(slide_lines),
                    )
                )
        return ParsedDocument(file_path=file_path, doc_type="pptx", segments=segments)

    if suffix == ".ppt":
        try:
            # Attempt to parse .ppt (97-2003) files with python-pptx
            prs = Presentation(str(file_path))
            segments: list[TextSegment] = []
            for i, slide in enumerate(prs.slides):
                slide_lines: list[str] = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        slide_lines.append(text.strip())
                if slide_lines:
                    segments.append(
                        TextSegment(
                            index=i,
                            source_unit=f"slide_{i+1}",
                            text="\n".join(slide_lines),
                        )
                    )
            return ParsedDocument(file_path=file_path, doc_type="ppt", segments=segments)
        except Exception:  # noqa: BLE001
            # If parsing fails (format not supported), return empty segments
            return ParsedDocument(file_path=file_path, doc_type="ppt", segments=[])

    return ParsedDocument(file_path=file_path, doc_type=suffix.lstrip("."), segments=[])


def extract_text(file_path: Path) -> str:
    return parse_document(file_path).full_text


def docx_to_html(file_path: Path) -> str:
    """将DOCX文档转换为HTML格式"""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from html import escape
        
        doc = Document(str(file_path))
        html_parts = ['<div style="font-family: Arial, sans-serif; line-height: 1.6; color: #333;">']
        
        for para in doc.paragraphs:
            if not para.text.strip():
                html_parts.append('<p>&nbsp;</p>')
                continue
            
            # 获取段落样式信息
            style = para.style
            style_name = style.name if style else "Normal"
            
            # 根据样式添加不同的HTML标签
            if style_name.startswith("Heading"):
                level = 1
                if "Heading" in style_name:
                    try:
                        level = int(style_name[-1])
                    except ValueError:
                        level = 1
                html_parts.append(f'<h{level} style="margin: 12px 0; font-weight: bold;">{escape(para.text)}</h{level}>')
            else:
                html_parts.append(f'<p style="margin: 8px 0;">{escape(para.text)}</p>')
        
        # 处理表格
        for table in doc.tables:
            html_parts.append('<table style="border-collapse: collapse; width: 100%; margin: 12px 0;">')
            for row in table.rows:
                html_parts.append('<tr>')
                for cell in row.cells:
                    html_parts.append(f'<td style="border: 1px solid #ccc; padding: 8px;">{escape(cell.text)}</td>')
                html_parts.append('</tr>')
            html_parts.append('</table>')
        
        html_parts.append('</div>')
        return "\n".join(html_parts)
    except Exception as e:
        return f'<div style="color: red; padding: 20px;">无法转换DOCX文件: {str(e)}</div>'



def extract_pdf_text_chunked(file_path: Path, max_pdf_pages: int = 20) -> dict:
    """
    提取PDF文本，分页面返回
    
    返回:
    {
        "total_pages": int,
        "extracted_pages": [{"page_num": int, "text": str, "char_count": int}, ...]
    }
    """
    try:
        reader = PdfReader(str(file_path), strict=False)
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                return {"total_pages": 0, "extracted_pages": []}
        
        total_pages = len(reader.pages)
        pages_to_extract = min(max_pdf_pages, total_pages)
        
        extracted_pages = []
        for i in range(pages_to_extract):
            page = reader.pages[i]
            page_text = (page.extract_text() or "").strip()
            if page_text:
                extracted_pages.append({
                    "page_num": i + 1,
                    "text": page_text,
                    "char_count": len(page_text)
                })
        
        return {
            "total_pages": total_pages,
            "extracted_pages": extracted_pages,
            "extracted_count": len(extracted_pages)
        }
    except Exception as e:
        logging.error(f"PDF text extraction failed: {str(e)}")
        return {"total_pages": 0, "extracted_pages": [], "error": str(e)}

# new function to analyze PDF with LLM and extract key insights for teachers
def analyze_pdf_with_llm(pdf_text: str, llm_client) -> dict:
    """
    使用LLM分析PDF文本，提取关键内容、要点、总结
    
    返回:
    {
        "status": "success|error",
        "summary": str,  # 整体总结
        "key_points": [str],  # 关键要点列表
        "focus_areas": [str],  # 重点领域
        "insights": str,  # 深度见解
        "analysis_time": float
    }
    """
    if not llm_client or not llm_client.enabled:
        return {
            "status": "skipped",
            "message": "LLM未启用"
        }
    
    import time
    start_time = time.time()
    
    try:
        # 限制文本长度用于分析
        analysis_text = pdf_text[:5000] if len(pdf_text) > 5000 else pdf_text
        
        # 第一步：生成总结
        system_prompt = "你是一名专业的文档分析师。你的任务是理解并分析文档内容，提取关键信息。"
        
        user_prompt = f"""请分析以下PDF文档内容，提供：
1. 一句话总结（不超过30字）
2. 3-5个关键要点（每个不超过20字）
3. 重点关注的领域（最多5个）
4. 对教师有价值的深度见解（不超过100字）

文档内容：
{analysis_text}

请返回JSON格式：
{{
    "summary": "总结文本",
    "key_points": ["点1", "点2", "点3"],
    "focus_areas": ["领域1", "领域2"],
    "insights": "深度见解"
}}"""
        
        response = llm_client.chat_json(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            model=None,
            temperature=0.3
        )
        
        analysis_time = time.time() - start_time
        
        if response:
            return {
                "status": "success",
                "summary": response.get("summary", ""),
                "key_points": response.get("key_points", []),
                "focus_areas": response.get("focus_areas", []),
                "insights": response.get("insights", ""),
                "analysis_time": analysis_time
            }
        else:
            return {
                "status": "error",
                "message": "LLM返回空响应",
                "analysis_time": analysis_time
            }
            
    except Exception as e:
        logging.error(f"PDF LLM analysis failed: {str(e)}")
        return {
            "status": "error",
            "message": str(e),
            "analysis_time": time.time() - start_time
        }


def get_pdf_key_insights(file_path: Path, llm_client) -> dict:
    """
    获取PDF的关键见解和摘要
    
    返回:
    {
        "status": "success|error",
        "pdf_stats": {...},
        "text_content": str,
        "analysis": {...}
    }
    """
    try:
        # 1. 提取PDF文本
        text_chunks = extract_pdf_text_chunked(file_path, max_pdf_pages=10)
        if not text_chunks.get("extracted_pages"):
            return {
                "status": "error",
                "message": "无法提取PDF文本内容"
            }
        
        # 合并所有页面文本
        full_text = "\n".join([p["text"] for p in text_chunks["extracted_pages"]])
        
        # 2. 使用LLM分析
        analysis = analyze_pdf_with_llm(full_text, llm_client)
        
        return {
            "status": "success",
            "pdf_stats": {
                "total_pages": text_chunks["total_pages"],
                "extracted_pages": text_chunks["extracted_count"],
                "total_chars": len(full_text)
            },
            "text_content": full_text,
            "analysis": analysis
        }
        
    except Exception as e:
        logging.error(f"PDF key insights extraction failed: {str(e)}")
        return {
            "status": "error",
            "message": str(e)
        }


def get_file_preview_data(file_path: Path, file_type: str) -> dict:
    """生成不同格式文件的预览数据"""
    result = {
        "type": file_type,
        "status": "success",
    }
    
    try:
        if file_type == "pdf":
            # PDF：返回base64编码的PDF内容和文件大小
            with open(file_path, "rb") as f:
                pdf_content = f.read()
            result["pdf_base64"] = base64.b64encode(pdf_content).decode("utf-8")
            result["file_size"] = len(pdf_content)
            
            # 获取页数
            reader = PdfReader(str(file_path), strict=False)
            result["page_count"] = len(reader.pages)
            
        elif file_type == "docx":
            # DOCX：转换为HTML
            html_content = docx_to_html(file_path)
            result["html_content"] = html_content
            result["file_size"] = file_path.stat().st_size
            
        elif file_type in ["pptx", "ppt"]:
            # PPT：转换为HTML（每页一个div）
            try:
                prs = Presentation(str(file_path))
                slides_html = []
                
                for slide_idx, slide in enumerate(prs.slides):
                    slide_content = f'<div style="page-break-after: always; padding: 20px; border: 1px solid #ddd; margin-bottom: 20px; background: #f9f9f9;">'
                    slide_content += f'<h2 style="margin-top: 0; color: #2196f3;">幻灯片 {slide_idx + 1}</h2>'
                    
                    slide_text = []
                    for shape in slide.shapes:
                        text = getattr(shape, "text", "").strip()
                        if text:
                            slide_text.append(f'<p style="margin: 8px 0;">{escape(text)}</p>')
                    
                    if slide_text:
                        slide_content += "\n".join(slide_text)
                    else:
                        slide_content += '<p style="color: #999;">(该页无文本内容)</p>'
                    
                    slide_content += '</div>'
                    slides_html.append(slide_content)
                
                result["html_content"] = "\n".join(slides_html)
                result["slide_count"] = len(prs.slides)
                result["file_size"] = file_path.stat().st_size
                
            except Exception as e:
                result["status"] = "partial"
                result["error"] = str(e)
                result["html_content"] = f'<div style="color: red; padding: 20px;">无法转换PPT文件: {str(e)}</div>'
        
        else:
            result["status"] = "unsupported"
            result["message"] = f"不支持的文件格式: {file_type}"
        
    except Exception as e:
        result["status"] = "error"
        result["error"] = str(e)
    
    return result
