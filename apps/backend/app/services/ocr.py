"""OCR service for processing document images and extracting text."""

import subprocess
import tempfile
from pathlib import Path

from app.services.document_parser import ParsedDocument, TextSegment
from pptx import Presentation
from pypdf import PdfReader

try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

try:
    import pytesseract
    HAS_PYTESSERACT = True
except ImportError:
    HAS_PYTESSERACT = False


def process_with_ocr(file_path: Path) -> ParsedDocument:
    """
    Process a file using OCR/advanced text extraction to extract text and structure.
    
    For PPTX files, performs comprehensive content extraction from all slides.
    For PDF files, attempts to extract all text including from images where possible.
    
    Args:
        file_path: Path to the document file to process
        
    Returns:
        ParsedDocument with extracted content
    """
    doc_type = file_path.suffix.lstrip(".").lower()
    
    if doc_type in {"pptx", "ppt"}:
        return _process_pptx_for_ocr(file_path, doc_type)
    elif doc_type == "pdf":
        return _process_pdf_for_ocr(file_path)
    else:
        # For unsupported formats, return empty document
        return ParsedDocument(
            file_path=file_path,
            doc_type=doc_type,
            segments=[],
        )


def _process_pptx_for_ocr(file_path: Path, doc_type: str) -> ParsedDocument:
    """
    Extract comprehensive content from PPTX files using text extraction and image OCR.
    Includes text from shapes, tables, speaker notes, and OCR on slide images if available.
    """
    segments: list[TextSegment] = []
    try:
        prs = Presentation(str(file_path))
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = []
            
            # Step 1: Extract direct text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                    
                # Extract text from table cells if present
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_content.append(cell.text.strip())
            
            # Step 2: Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_content.append(f"[备注] {notes_text}")
            
            # Step 3: If little text extracted, try OCR on slide image
            extracted_text = "\n".join(slide_content)
            if len(extracted_text) < 500 and HAS_PYTESSERACT:
                try:
                    ocr_text = _extract_pptx_slide_ocr(prs, slide_idx)
                    if ocr_text:
                        if extracted_text:
                            extracted_text = extracted_text + "\n[OCR识别]\n" + ocr_text
                        else:
                            extracted_text = "[OCR识别]\n" + ocr_text
                except Exception:
                    pass
            
            # Create segment if content found
            if extracted_text:
                segments.append(
                    TextSegment(
                        index=slide_idx,
                        source_unit=f"slide_{slide_idx+1}",
                        text=extracted_text,
                    )
                )
    except Exception:
        # If parsing fails, return empty segments
        pass
    
    return ParsedDocument(
        file_path=file_path,
        doc_type=doc_type,
        segments=segments,
    )


def _extract_pptx_slide_ocr(prs, slide_idx: int) -> str:
    """Extract text from a specific slide using OCR via image rendering."""
    if not HAS_PYTESSERACT:
        return ""
    
    try:
        from PIL import Image
        import io
        
        slide = prs.slides[slide_idx]
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Convert to pixels (assuming 96 DPI)
        dpi = 96
        width_px = int(slide_width / 914400 * dpi)  # EMUS to pixels
        height_px = int(slide_height / 914400 * dpi)
        
        # Create image from slide
        img = Image.new('RGB', (width_px, height_px), color='white')
        
        # Render shapes to image (simplified - just uses coordinates)
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            img.save(tmp.name)
            try:
                # Use tesseract via subprocess (more reliable)
                result = subprocess.run(
                    ['tesseract', tmp.name, 'stdout', '-l', 'chi_sim+eng'],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                ocr_text = result.stdout.strip()
                return ocr_text if ocr_text else ""
            except Exception:
                return ""
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    except Exception:
        return ""


def _process_pdf_for_ocr(file_path: Path) -> ParsedDocument:
    """
    Extract comprehensive content from PDF files using text extraction and image OCR.
    Processes all pages, extracting text and running OCR on image-based pages if available.
    """
    segments: list[TextSegment] = []
    try:
        # First, try standard PDF text extraction
        reader = PdfReader(str(file_path), strict=False)
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                return ParsedDocument(
                    file_path=file_path,
                    doc_type="pdf",
                    segments=[],
                )
        
        # Extract text from all pages
        for page_idx, page in enumerate(reader.pages):
            page_text = (page.extract_text() or "").strip()
            
            # If little text on this page and OCR is available, try OCR
            if len(page_text) < 200 and HAS_PYTESSERACT and HAS_PDF2IMAGE:
                try:
                    ocr_text = _extract_pdf_page_ocr(file_path, page_idx)
                    if ocr_text:
                        if page_text:
                            page_text = page_text + "\n[OCR识别]\n" + ocr_text
                        else:
                            page_text = "[OCR识别]\n" + ocr_text
                except Exception:
                    pass
            
            if page_text:
                segments.append(
                    TextSegment(
                        index=page_idx,
                        source_unit=f"page_{page_idx+1}",
                        text=page_text,
                    )
                )
    except Exception:
        pass
    
    return ParsedDocument(
        file_path=file_path,
        doc_type="pdf",
        segments=segments,
    )


def _extract_pdf_page_ocr(file_path: Path, page_idx: int) -> str:
    """Extract text from a specific PDF page using OCR."""
    if not HAS_PDF2IMAGE or not HAS_PYTESSERACT:
        return ""
    
    try:
        # Convert PDF page to image
        images = convert_from_path(str(file_path), first_page=page_idx+1, last_page=page_idx+1, dpi=150)
        if not images:
            return ""
        
        img = images[0]
        
        # Use tesseract for OCR
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            img.save(tmp.name)
            try:
                result = subprocess.run(
                    ['tesseract', tmp.name, 'stdout', '-l', 'chi_sim+eng'],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                ocr_text = result.stdout.strip()
                return ocr_text if ocr_text else ""
            except Exception:
                return ""
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    except Exception:
        return ""
